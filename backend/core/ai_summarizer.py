"""
AI Summarizer Module
Summarizes content from various inspiration types using AI
"""

import base64
import os
import json
from pathlib import Path
from typing import Optional, Dict, Any, List
from abc import ABC, abstractmethod

from ..models import Inspiration, InspirationType, AIModelConfig


TEXT_MODE_TYPES = [
    InspirationType.CODE, InspirationType.TEXT, InspirationType.ENVIRONMENT,
    InspirationType.DATA, InspirationType.DOCUMENT, InspirationType.NOTEBOOK,
    InspirationType.SCRIPT, InspirationType.STYLE, InspirationType.MARKUP,
    InspirationType.FOLDER, InspirationType.CONFIG
]


class BaseSummarizer(ABC):
    @abstractmethod
    async def summarize(self, content: Any, **kwargs) -> str:
        pass


class ImageSummarizer(BaseSummarizer):
    def __init__(self, config: AIModelConfig):
        self.config = config
    
    def _encode_image(self, image_path: str) -> str:
        with open(image_path, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    
    async def summarize(self, image_path: str, **kwargs) -> str:
        from openai import AsyncOpenAI
        
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url
        )
        
        # Check if model config is valid
        if not self.config.api_key:
             return json.dumps({
                "tree": "",
                "overview": "错误：未配置 API Key。请在设置中配置 AI 模型。",
                "important_docs": "",
                "secondary_docs": "",
                "_context": {}
            }, ensure_ascii=False)
        
        base64_image = self._encode_image(image_path)
        
        response = await client.chat.completions.create(
            model=self.config.model_name,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请详细描述这张图片的内容，包括：1. 主要元素和主题 2. 色彩和构图特点 3. 可能传达的情感或意境 4. 适合用于什么类型的创意项目"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000
        )
        
        return response.choices[0].message.content


class DocumentSummarizer(BaseSummarizer):
    def __init__(self, config: AIModelConfig):
        self.config = config
    
    def _read_docx(self, file_path: str) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        content.append(row_text)
            return '\n'.join(content)
        except Exception as e:
            return f"[无法解析Word文档: {str(e)}]"
    
    def _read_pdf(self, file_path: str) -> str:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            content = []
            for page in reader.pages:
                text = page.extract_text()
                if text and text.strip():
                    content.append(text)
            return '\n'.join(content)
        except Exception as e:
            return f"[无法解析PDF文档: {str(e)}]"
    
    def _read_content(self, file_path: str) -> str:
        path = Path(file_path)
        ext = path.suffix.lower()
        
        if ext == '.docx':
            return self._read_docx(file_path)
        elif ext == '.pdf':
            return self._read_pdf(file_path)
        elif path.is_file():
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception:
                try:
                    with open(path, 'rb') as f:
                        content = f.read()
                        try:
                            return content.decode('utf-8')
                        except UnicodeDecodeError:
                            return content.decode('latin-1', errors='ignore')
                except Exception:
                    return ""
        elif path.is_dir():
            content_parts = []
            
            tree_structure = self._get_folder_structure(path)
            content_parts.append(f"=== 文件夹结构 ===\n{tree_structure}\n")
            
            for ext in ['.txt', '.md', '.py', '.js', '.ts', '.json', '.yaml', '.yml', '.xml', '.html', '.css', '.vue', '.jsx', '.tsx', '.java', '.go', '.rs', '.c', '.cpp', '.h', '.hpp', '.sh', '.bat', '.ps1', '.env', '.ini', '.cfg', '.toml']:
                for file_path in path.rglob(f'*{ext}'):
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            relative_path = file_path.relative_to(path)
                            content_parts.append(f"=== {relative_path} ===\n{f.read()[:3000]}")
                    except Exception:
                        continue
            
            return "\n\n".join(content_parts[:15])
    
    def _get_folder_structure(self, folder_path: Path, prefix: str = "", max_depth: int = 3, current_depth: int = 0) -> str:
        if current_depth >= max_depth:
            return ""
        
        structure_lines = []
        try:
            items = sorted(folder_path.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower()))
            for item in items[:50]:
                if item.name.startswith('.') or item.name in ['node_modules', '__pycache__', '.git', 'venv', 'build', 'dist', '.venv']:
                    continue
                indent = "  " * current_depth
                if item.is_dir():
                    structure_lines.append(f"{indent}+ {item.name}/")
                    sub_structure = self._get_folder_structure(item, prefix, max_depth, current_depth + 1)
                    if sub_structure:
                        structure_lines.append(sub_structure)
                else:
                    size = item.stat().st_size
                    if size < 1024:
                        size_str = f"{size}B"
                    elif size < 1024 * 1024:
                        size_str = f"{size / 1024:.1f}KB"
                    else:
                        size_str = f"{size / (1024 * 1024):.1f}MB"
                    structure_lines.append(f"{indent}- {item.name} ({size_str})")
        except PermissionError:
            pass
        
        return "\n".join(structure_lines)
    
    async def summarize(self, file_path: str, file_type: str = "document", **kwargs) -> str:
        from openai import AsyncOpenAI
        
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url
        )
        
        content = self._read_content(file_path)
        
        if not content:
            return "无法读取文件内容"
        
        if len(content) > 10000:
            content = content[:10000] + "\n... (内容已截断)"
        
        type_prompts = {
            "code": "你是一个代码分析专家，请分析代码并总结其功能、结构和特点。",
            "text": "你是一个文本分析专家，请提取文本的关键信息。",
            "document": "你是一个文档分析专家，请分析文档内容并提取关键信息。",
            "notebook": "你是一个Jupyter笔记本分析专家，请分析笔记本内容并总结。",
            "script": "你是一个脚本分析专家，请分析脚本功能和使用方法。",
            "style": "你是一个样式分析专家，请分析CSS/样式表的结构和设计特点。",
            "markup": "你是一个标记语言分析专家，请分析HTML/XML结构和内容。",
            "data": "你是一个数据分析专家，请分析数据文件的结构和内容。",
            "environment": "你是一个配置分析专家，请分析配置文件的内容和用途。",
            "config": "你是一个配置分析专家，请分析配置文件的内容和用途。",
            "folder": """你是一个项目分析专家，请分析文件夹内容并生成详细的项目结构总结。

请按照以下格式输出：

## 项目概览
[简要描述这个项目/文件夹的用途和主题]

## 目录结构
[分析目录结构的特点，如模块划分、层次结构等]

## 主要文件分析
[分析关键文件的功能和作用]

## 技术栈/文件类型
[识别使用的技术、语言或文件类型]

## 创意灵感点
[提取可用于创意生成的灵感点和特色]

## 潜在应用场景
[这个项目/素材可以用于什么样的创意项目]"""
        }
        
        system_prompt = type_prompts.get(file_type, "你是一个内容分析专家，请分析并总结内容。")
        
        user_prompt = f"请分析以下内容并总结：\n\n1. 主要内容和主题\n2. 关键信息和要点\n3. 结构和特点\n4. 可用于创意的灵感点\n\n内容：\n{content}"
        
        if file_type == "folder":
            user_prompt = f"请分析以下文件夹内容并生成详细的项目结构总结：\n\n{content}"
        
        max_tokens_value = 10000 if file_type == "folder" else 2000
        
        response = await client.chat.completions.create(
            model=self.config.model_name,
            messages=[
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": user_prompt
                }
            ],
            max_tokens=max_tokens_value
        )
        
        return response.choices[0].message.content


class FolderSummarizer(BaseSummarizer):
    def __init__(self, config: AIModelConfig):
        self.config = config
    
    def _get_folder_tree(self, folder_path: Path, prefix: str = "", max_depth: int = 10, current_depth: int = 0, ignored_paths: List[str] = None) -> str:
        if current_depth >= max_depth:
            return ""
        
        ignored_paths = ignored_paths or []
        
        def is_ignored(item_path: Path) -> bool:
            try:
                relative = str(item_path.relative_to(folder_path))
                for ignored in ignored_paths:
                    if relative == ignored or relative.startswith(ignored + '/'):
                        return True
            except ValueError:
                pass
            return False
        
        structure_lines = []
        try:
            items = sorted(folder_path.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower()))
            for item in items:
                if item.name.startswith('.') or item.name in ['node_modules', '__pycache__', '.git', 'venv', 'build', 'dist', '.venv', '__init__.py', 'package-lock.json', 'yarn.lock']:
                    continue
                if is_ignored(item):
                    continue
                indent = "  " * current_depth
                if item.is_dir():
                    structure_lines.append(f"{indent}+ {item.name}/")
                    sub_structure = self._get_folder_tree(item, prefix, max_depth, current_depth + 1, ignored_paths)
                    if sub_structure:
                        structure_lines.append(sub_structure)
                else:
                    structure_lines.append(f"{indent}- {item.name}")
        except PermissionError:
            pass
        
        return "\n".join(structure_lines)
    
    def _get_file_info(self, file_path: Path) -> Dict[str, Any]:
        try:
            stat = file_path.stat()
            size = stat.st_size
            size_str = f"{size}B" if size < 1024 else f"{size/1024:.1f}KB" if size < 1024*1024 else f"{size/(1024*1024):.1f}MB"
            return {
                "size": size,
                "size_str": size_str,
                "extension": file_path.suffix.lower(),
                "modified": stat.st_mtime
            }
        except Exception:
            return {"size": 0, "size_str": "unknown", "extension": file_path.suffix.lower(), "modified": 0}
    
    def _read_file_content(self, file_path: Path, max_length: int = 6000) -> str:
        try:
            suffix = file_path.suffix.lower()
            
            if suffix == '.docx':
                try:
                    from docx import Document
                    doc = Document(str(file_path))
                    content = []
                    for para in doc.paragraphs[:30]:
                        if para.text.strip():
                            content.append(para.text)
                    return '\n'.join(content)[:max_length]
                except Exception as e:
                    return f"[Word文档解析失败: {str(e)}]"
            
            elif suffix == '.pdf':
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(str(file_path))
                    content = []
                    for page in reader.pages[:10]:
                        text = page.extract_text()
                        if text and text.strip():
                            content.append(text)
                    return '\n'.join(content)[:max_length]
                except Exception as e:
                    return f"[PDF解析失败: {str(e)}]"
            
            elif suffix in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(str(file_path), read_only=True)
                    content = []
                    for sheet_name in wb.sheetnames[:3]:
                        sheet = wb[sheet_name]
                        content.append(f"=== Sheet: {sheet_name} ===")
                        for row in list(sheet.iter_rows(max_row=20, values_only=True)):
                            if any(cell for cell in row if cell):
                                content.append(' | '.join(str(cell) if cell else '' for cell in row))
                    return '\n'.join(content)[:max_length]
                except Exception as e:
                    return f"[Excel解析失败: {str(e)}]"
            
            elif suffix in ['.pptx', '.ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(str(file_path))
                    content = []
                    for i, slide in enumerate(prs.slides[:10]):
                        content.append(f"=== Slide {i+1} ===")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                content.append(shape.text)
                    return '\n'.join(content)[:max_length]
                except Exception as e:
                    return f"[PPT解析失败: {str(e)}]"
            
            elif suffix in ['.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a']:
                info = self._get_file_info(file_path)
                return f"[音频文件: {info['size_str']}, 格式: {suffix}]"
            
            elif suffix in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm']:
                info = self._get_file_info(file_path)
                return f"[视频文件: {info['size_str']}, 格式: {suffix}]"
            
            elif suffix in ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg', '.ico', '.tiff']:
                info = self._get_file_info(file_path)
                return f"[图片文件: {info['size_str']}, 格式: {suffix}]"
            
            elif suffix in ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']:
                info = self._get_file_info(file_path)
                return f"[压缩文件: {info['size_str']}, 格式: {suffix}]"
            
            elif suffix in ['.pt', '.pth', '.h5', '.hdf5', '.onnx', '.pb', '.safetensors', '.bin', '.ckpt', '.pkl']:
                info = self._get_file_info(file_path)
                return f"[模型文件: {info['size_str']}, 格式: {suffix}]"
            
            elif suffix in ['.exe', '.dll', '.so', '.dylib', '.bin']:
                info = self._get_file_info(file_path)
                return f"[二进制文件: {info['size_str']}, 格式: {suffix}]"
            
            else:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                        if not content.strip():
                            info = self._get_file_info(file_path)
                            return f"[空文件或二进制: {info['size_str']}]"
                        return content[:max_length]
                except Exception:
                    info = self._get_file_info(file_path)
                    return f"[文件读取失败: {info['size_str']}, 格式: {suffix}]"
                    
        except Exception as e:
            return f"[文件处理错误: {str(e)}]"
    
    def _get_file_importance(self, file_path: Path) -> str:
        name = file_path.name.lower()
        parent = file_path.parent.name.lower() if file_path.parent else ""
        
        important_names = ['readme', 'main', 'index', 'app', 'config', 'setup', 
                          'requirements', 'package', 'cargo', 'go.mod', 'pom',
                          'build', 'makefile', 'dockerfile', 'docker-compose']
        important_dirs = ['src', 'lib', 'core', 'main', 'app', 'api']
        
        if any(imp in name for imp in important_names):
            return "important"
        if parent in important_dirs:
            return "important"
        if name.startswith('test') or name.startswith('spec'):
            return "secondary"
        if file_path.suffix.lower() in ['.md', '.txt', '.rst', '.doc', '.docx', '.pdf']:
            return "important"
        
        return "normal"
    
    async def _summarize_file(self, client, file_path: Path, relative_path: str, content: str) -> str:
        try:
            suffix = file_path.suffix.lower()
            
            if content.startswith('[') and content.endswith(']'):
                prompt = f"""文件: {relative_path}
文件信息: {content}

请简要描述这个文件的类型和用途（1-2句话）。"""
            else:
                prompt = f"""文件: {relative_path}

内容:
{content[:3000]}

请分析这个文件的内容并总结其功能和用途（2-3句话）。"""
            
            response = await client.chat.completions.create(
                model=self.config.model_name,
                messages=[
                    {
                        "role": "system",
                        "content": "你是一个文件分析专家，请简要分析文件内容并总结其功能和用途。"
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                max_tokens=300
            )
            return response.choices[0].message.content
        except Exception as e:
            print(f"Error summarizing file {relative_path}: {e}")
            return f"[总结失败: {str(e)}]"
    
    async def _summarize_recursive(self, current_path: Path, root_path: Path, client, semaphore, ignored_paths: List[str]) -> Dict[str, Any]:
        import asyncio
        
        try:
            if current_path == root_path:
                relative_path = ""
            else:
                relative_path = str(current_path.relative_to(root_path)).replace('\\', '/')
        except ValueError:
            relative_path = current_path.name
        
        def is_ignored(path_to_check: str) -> bool:
            if not path_to_check:
                return False
            for ignored in ignored_paths:
                if path_to_check == ignored or path_to_check.startswith(ignored + '/'):
                    return True
            return False
        
        if is_ignored(relative_path):
            print(f"Skipping ignored path: {relative_path}")
            return None

        if current_path.is_file():
            if current_path.name.startswith('.') or current_path.name in ['package-lock.json', 'yarn.lock', '.DS_Store', 'Thumbs.db']:
                print(f"DEBUG: Skipping system file {current_path.name}")
                return None
                
            async with semaphore:
                try:
                    print(f"DEBUG: Processing file: {relative_path}")
                    content = self._read_file_content(current_path, max_length=6000)
                    
                    summary = await self._summarize_file(client, current_path, relative_path, content)
                    
                    print(f"DEBUG: Successfully summarized file: {relative_path}")
                    return {
                        "path": relative_path,
                        "name": current_path.name,
                        "type": "file",
                        "summary": summary,
                        "importance": self._get_file_importance(current_path)
                    }
                except Exception as e:
                    print(f"DEBUG: Error summarizing file {current_path}: {e}")
                    return None

        if current_path.is_dir():
            print(f"DEBUG: Processing directory: {relative_path if relative_path else 'ROOT'}")
            
            direct_children = []
            try:
                items = sorted(current_path.iterdir(), key=lambda x: (not x.is_dir(), x.name.lower()))
                print(f"DEBUG: Found {len(items)} items in {current_path.name}")
                
                for item in items:
                    if item.name.startswith('.') or item.name in ['node_modules', '__pycache__', '.git', 'dist', 'build', 'venv', '.venv']:
                        print(f"DEBUG: Skipping excluded item: {item.name}")
                        continue
                    
                    try:
                        if current_path == root_path:
                            item_relative = item.name
                        else:
                            item_relative = f"{relative_path}/{item.name}"
                    except Exception:
                        item_relative = item.name
                    
                    if is_ignored(item_relative):
                        print(f"DEBUG: Skipping ignored item: {item_relative}")
                        continue
                    
                    direct_children.append(item)
                
                if not direct_children:
                    print(f"DEBUG: No valid children for directory: {relative_path}")
                    return {
                        "path": relative_path,
                        "name": current_path.name,
                        "type": "folder",
                        "summary": "空文件夹",
                        "children": []
                    }

                child_tasks = [
                    self._summarize_recursive(item, root_path, client, semaphore, ignored_paths)
                    for item in direct_children
                ]
                
                results = await asyncio.gather(*child_tasks)
                children = [r for r in results if r is not None]
                print(f"DEBUG: Got {len(children)} valid children for {current_path.name}")
                
                if not children:
                    print(f"DEBUG: No children summaries for directory: {relative_path}")
                    if current_path == root_path:
                        return {
                            "path": relative_path,
                            "name": current_path.name,
                            "type": "folder",
                            "summary": "文件夹为空或所有内容被忽略",
                            "children": []
                        }
                    return None
                
                direct_children_summaries = []
                for child in children:
                    child_summary = child.get('summary', '')[:150]
                    direct_children_summaries.append(
                        f"- [{child['type'].upper()}] {child['name']}: {child_summary}"
                    )
                
                children_info = "\n".join(direct_children_summaries)
                
                if len(children_info) > 8000:
                    children_info = children_info[:8000] + "\n...(内容已截断)"
                
                async with semaphore:
                    prompt = f"""请总结以下文件夹的内容。

文件夹路径: {relative_path if relative_path else "Root"}

该文件夹包含以下直接子项（文件和子文件夹）:
{children_info}

请生成一个简洁的文件夹总结（150-200字），概括：
1. 该文件夹的主要用途
2. 包含的主要内容类型
3. 各子项之间的关系（如果有）"""

                    try:
                        print(f"Generating summary for directory: {relative_path if relative_path else 'ROOT'}")
                        response = await client.chat.completions.create(
                            model=self.config.model_name,
                            messages=[
                                {"role": "system", "content": "你是一个项目架构分析专家，擅长总结文件夹结构和内容。"},
                                {"role": "user", "content": prompt}
                            ],
                            max_tokens=500
                        )
                        folder_summary = response.choices[0].message.content
                    except Exception as e:
                        print(f"Error generating folder summary for {relative_path}: {e}")
                        folder_summary = f"总结生成失败: {str(e)}"

                return {
                    "path": relative_path,
                    "name": current_path.name,
                    "type": "folder",
                    "summary": folder_summary,
                    "children": children
                }
            except Exception as e:
                print(f"Error processing directory {current_path}: {e}")
                return None
        
        return None

    async def summarize(self, folder_path: str, ignored_paths: List[str] = None, **kwargs) -> str:
        from openai import AsyncOpenAI
        import asyncio
        
        # Check if model config is valid
        if not self.config.api_key:
            return json.dumps({
                "tree": "",
                "overview": "错误：未配置 API Key。请在设置中配置 AI 模型。",
                "important_docs": "",
                "secondary_docs": "",
                "_context": {}
            }, ensure_ascii=False)
        
        # Increase timeout to avoid timeouts during long processing
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url,
            timeout=120.0
        )
        
        path = Path(folder_path)
        if not path.is_dir():
            return json.dumps({
                "tree": "",
                "overview": f"错误：路径不是有效的文件夹 ({folder_path})",
                "important_docs": "",
                "secondary_docs": "",
                "_context": {}
            }, ensure_ascii=False)
        
        ignored_paths = ignored_paths or []
        
        # Semaphore for concurrency control
        sem = asyncio.Semaphore(10) # Increase concurrency slightly as we are doing hierarchical
        
        # Recursive summarization
        root_summary = await self._summarize_recursive(path, path, client, sem, ignored_paths)
        
        if not root_summary:
            return json.dumps({
                "tree": "",
                "overview": "无法生成总结（可能是空文件夹或所有文件被忽略）。",
                "important_docs": "",
                "secondary_docs": "",
                "_context": {}
            }, ensure_ascii=False)

        # Flatten results for file_summaries (frontend compatibility)
        # We include both files and folders in the flat list so frontend can display them
        flat_summaries = []
        
        def flatten(node):
            if not node: return
            flat_summaries.append({
                "path": node['path'],
                "name": node['name'],
                "summary": node['summary'],
                "type": node['type']
            })
            if node.get('children'):
                for child in node['children']:
                    flatten(child)
        
        flatten(root_summary)
        
        # Generate Tree (standard visual tree)
        tree_structure = self._get_folder_tree(path, ignored_paths=ignored_paths)
        
        # Overview is the root folder summary
        overview = root_summary.get('summary', '')
        
        # Construct Result - Simplified as requested
        result = {
            "tree": tree_structure,
            "overview": overview,
            "important_docs": "", # Removed as requested
            "secondary_docs": "", # Removed as requested
            "_context": {
                "file_summaries": flat_summaries,
            }
        }
        
        return json.dumps(result, ensure_ascii=False)

    async def regenerate_section(self, section: str, context: Dict[str, Any], folder_path: str = None, ignored_paths: List[str] = None) -> str:
        from openai import AsyncOpenAI
        
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url,
            timeout=120.0
        )
        
        file_summaries = context.get('file_summaries', [])
        tree_structure = context.get('tree', '') # context might have tree if passed fully
        
        if section == 'tree':
            if folder_path:
                return self._get_folder_tree(Path(folder_path), ignored_paths=ignored_paths)
            return tree_structure
        elif section == 'overview':
            prompt = f"""请根据以下文件夹结构和文件摘要，重新生成一份项目概览。
            
文件夹结构：
{tree_structure}

文件摘要：
{json.dumps([{'path': fs['path'], 'summary': fs['summary']} for fs in file_summaries], ensure_ascii=False, indent=2)}

请生成一个“项目/素材概览”，包含：
1. 整体用途和主题
2. 目录结构特点分析
3. 核心功能或内容概述
"""
        elif section == 'important_docs':
            important_summaries = [fs for fs in file_summaries if fs.get('importance') == 'important']
            prompt = f"""请根据以下重要文件的摘要，重新生成一份“重要文档说明”。
            
重要文件摘要：
{json.dumps([{'path': fs['path'], 'summary': fs['summary']} for fs in important_summaries], ensure_ascii=False, indent=2)}

请详细说明这些核心文件的作用和它们之间的关系。
"""
        elif section == 'secondary_docs':
            secondary_summaries = [fs for fs in file_summaries if fs.get('importance') == 'secondary']
            prompt = f"""请根据以下次要文件的摘要，重新生成一份“次要文档说明”。
            
次要文件摘要：
{json.dumps([{'path': fs['path'], 'summary': fs['summary']} for fs in secondary_summaries], ensure_ascii=False, indent=2)}

请简要说明这些文件的作用（如配置、工具、日志等）。
"""
        else:
            return "无效的章节"
            
        response = await client.chat.completions.create(
            model=self.config.model_name,
            messages=[
                {"role": "system", "content": "你是一个项目分析专家。"},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000
        )
        
        return response.choices[0].message.content

    
    async def regenerate_all_summaries(self, folder_path: str, ignored_paths: List[str] = None) -> Dict[str, Any]:
        # Reuse summarize logic
        json_str = await self.summarize(folder_path, ignored_paths)
        try:
            data = json.loads(json_str)
            return {
                "overall_summary": data.get("overview", ""),
                "file_summaries": data.get("_context", {}).get("file_summaries", [])
            }
        except:
            return {
                "overall_summary": "生成失败",
                "file_summaries": []
            }
    
    async def regenerate_single_summary(self, folder_path: str, file_path: str, ignored_paths: List[str] = None) -> str:
        from openai import AsyncOpenAI
        import asyncio
        
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url,
            timeout=120.0
        )
        
        path = Path(folder_path)
        target_path = path / file_path if file_path else path
        
        if not target_path.exists():
            return "文件/文件夹不存在"
        
        ignored_paths = ignored_paths or []
        sem = asyncio.Semaphore(10)
        
        if target_path.is_file():
            try:
                content = self._read_file_content(target_path, max_length=6000)
                return await self._summarize_file(client, target_path, file_path, content)
            except Exception as e:
                return f"总结失败: {str(e)}"
         
        elif target_path.is_dir():
            result = await self._summarize_recursive(target_path, path, client, sem, ignored_paths)
            return result.get('summary', '') if result else "无法生成总结"
            
        return "未知类型"
    
    async def regenerate_node_summary(self, folder_path: str, node_path: str, ignored_paths: List[str] = None) -> Dict[str, Any]:
        from openai import AsyncOpenAI
        import asyncio
        
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url,
            timeout=120.0
        )
        
        path = Path(folder_path)
        target_path = path / node_path if node_path else path
        
        if not target_path.exists():
            return {"error": "文件/文件夹不存在"}
        
        ignored_paths = ignored_paths or []
        sem = asyncio.Semaphore(10)
        
        result = await self._summarize_recursive(target_path, path, client, sem, ignored_paths)
        
        if result:
            return {
                "path": result.get("path", ""),
                "name": result.get("name", ""),
                "type": result.get("type", ""),
                "summary": result.get("summary", ""),
                "children": result.get("children", [])
            }
        
        return {"error": "无法生成总结"}


class TextContentSummarizer(BaseSummarizer):
    def __init__(self, config: AIModelConfig):
        self.config = config
    
    def _read_content(self, file_path: str) -> str:
        path = Path(file_path)
        
        if path.is_file():
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception:
                try:
                    with open(path, 'rb') as f:
                        content = f.read()
                        try:
                            return content.decode('utf-8')
                        except UnicodeDecodeError:
                            return content.decode('latin-1', errors='ignore')
                except Exception:
                    return ""
        elif path.is_dir():
            content_parts = []
            
            tree_structure = self._get_folder_structure(path)
            content_parts.append(f"=== 文件夹结构 ===\n{tree_structure}\n")
            
            file_count = 0
            dir_count = 0
            file_type_counts = {}
            main_files = []
            
            for ext in ['.py', '.js', '.ts', '.json', '.yaml', '.yml', '.xml', '.html', '.css', '.vue', '.jsx', '.tsx', '.java', '.go', '.rs', '.c', '.cpp', '.h', '.hpp', '.sh', '.bat', '.ps1', '.env', '.ini', '.cfg', '.toml']:
                for file_path in path.rglob(f'*{ext}'):
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            relative_path = file_path.relative_to(path)
                            content = f.read()
                            if len(content) > 0:
                                file_count += 1
                                main_files.append(f"  - {relative_path}")
                    except Exception:
                        continue
            
            for ext in ['.txt', '.md', '.doc', '.docx', '.pdf', '.rst']:
                for file_path in path.rglob(f'*{ext}'):
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            relative_path = file_path.relative_to(path)
                            content = f.read()
                            if len(content) > 0:
                                file_count += 1
                                main_files.append(f"  - {relative_path}")
                    except Exception:
                        continue
            
            for item in path.iterdir():
                if item.is_dir() and not item.name.startswith('.') and item.name not in ['node_modules', '__pycache__', '.git', 'venv', 'build', 'dist', '.venv']:
                    dir_count += 1
                    try:
                        sub_items = list(item.iterdir())[:5]
                        sub_names = [sub.name for sub in sub_items if not sub.name.startswith('.')]
                        if sub_names:
                            content_parts.append(f"  - {item.name}/ ({', '.join(sub_names[:3])}{'...' if len(sub_names) > 3 else ''})")
                    except PermissionError:
                        pass
            
            if main_files:
                content_parts.append(f"\n主要文件 ({len(main_files)} 个):")
                for f in main_files[:10]:
                    content_parts.append(f"  {f}")
                if len(main_files) > 10:
                    content_parts.append(f"  ... 还有 {len(main_files) - 10} 个文件")
            
            content_parts.append(f"\n总计: {file_count} 个文件, {dir_count} 个文件夹")
            
            return "\n".join(content_parts[:20])
        
        return ""
    
    async def summarize(self, file_path: str, file_type: str = "text", **kwargs) -> str:
        from openai import AsyncOpenAI
        
        client = AsyncOpenAI(
            api_key=self.config.api_key,
            base_url=self.config.base_url
        )
        
        content = self._read_content(file_path)
        
        if not content:
            return "无法读取文件内容"
        
        if len(content) > 10000:
            content = content[:10000] + "\n... (内容已截断)"
        
        type_prompts = {
            "code": "你是一个代码分析专家，请分析代码并总结其功能、结构和特点。",
            "text": "你是一个文本分析专家，请提取文本的关键信息。",
            "document": "你是一个文档分析专家，请分析文档内容并提取关键信息。",
            "notebook": "你是一个Jupyter笔记本分析专家，请分析笔记本内容并总结。",
            "script": "你是一个脚本分析专家，请分析脚本功能和使用方法。",
            "style": "你是一个样式分析专家，请分析CSS/样式表的结构和设计特点。",
            "markup": "你是一个标记语言分析专家，请分析HTML/XML结构和内容。",
            "data": "你是一个数据分析专家，请分析数据文件的结构和内容。",
            "environment": "你是一个配置分析专家，请分析配置文件的内容和用途。",
            "config": "你是一个配置分析专家，请分析配置文件的内容和用途。",
            "folder": "你是一个项目分析专家，请分析文件夹内容并总结项目结构。"
        }
        
        system_prompt = type_prompts.get(file_type, "你是一个内容分析专家，请分析并总结内容。")
        
        response = await client.chat.completions.create(
            model=self.config.model_name,
            messages=[
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": f"请分析以下内容并总结：\n\n1. 主要内容和主题\n2. 关键信息和要点\n3. 结构和特点\n4. 可用于创意的灵感点\n\n内容：\n{content}"
                }
            ],
            max_tokens=2000
        )
        
        return response.choices[0].message.content


class AISummarizer:
    def __init__(self):
        self.model_configs: Dict[str, AIModelConfig] = {}
        self.summarizers: Dict[str, BaseSummarizer] = {}
        self.default_config: Optional[AIModelConfig] = None
        self.folder_summarizer: Optional[FolderSummarizer] = None
    
    def register_model(self, config: AIModelConfig):
        for file_type in config.file_types:
            self.model_configs[file_type] = config
            
            if file_type == "image":
                self.summarizers[file_type] = ImageSummarizer(config)
            elif file_type == "document":
                self.summarizers[file_type] = DocumentSummarizer(config)
            elif file_type == "folder":
                self.summarizers[file_type] = FolderSummarizer(config)
            elif file_type in ["code", "text", "notebook", "script", "style", "markup", "data", "environment", "config"]:
                self.summarizers[file_type] = TextContentSummarizer(config)
        
        if config.is_default:
            self.default_config = config
            self.folder_summarizer = FolderSummarizer(config)
    
    def get_summarizer(self, inspiration_type: str) -> Optional[BaseSummarizer]:
        if inspiration_type == "folder":
            # Prioritize specific folder summarizer if registered
            if "folder" in self.summarizers:
                return self.summarizers["folder"]
            # Fallback to default
            return self.folder_summarizer if self.folder_summarizer else FolderSummarizer(self.default_config) if self.default_config else None
        return self.summarizers.get(inspiration_type)
    
    async def summarize_inspiration(self, inspiration: Inspiration, ignored_paths: List[str] = None) -> str:
        summarizer = self.get_summarizer(inspiration.type)
        
        if not summarizer:
            if self.default_config:
                if inspiration.type == "folder":
                    summarizer = FolderSummarizer(self.default_config)
                else:
                    summarizer = DocumentSummarizer(self.default_config)
            else:
                return f"暂不支持 {inspiration.type} 类型的内容总结，请先配置AI模型"
        
        try:
            if isinstance(summarizer, FolderSummarizer):
                summary = await summarizer.summarize(inspiration.path, ignored_paths=ignored_paths)
            elif isinstance(summarizer, (TextContentSummarizer, DocumentSummarizer)):
                summary = await summarizer.summarize(inspiration.path, inspiration.type)
            else:
                summary = await summarizer.summarize(inspiration.path)
            return summary
        except Exception as e:
            return f"总结失败: {str(e)}"
    
    async def batch_summarize(self, inspirations: List[Inspiration]) -> Dict[str, str]:
        results = {}
        for inspiration in inspirations:
            ignored_paths = inspiration.metadata.get('ignored_paths', []) if inspiration.metadata else []
            results[inspiration.id] = await self.summarize_inspiration(inspiration, ignored_paths)
        return results
